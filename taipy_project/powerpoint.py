from pptx import Presentation
import PyPDF2
import fitz

def getInfo(file):
    if file.find(".pptx") >= 0:
        print("PPTX Read.")
        prs = Presentation(file)
        slides = prs.slides
        info = []

        for i in slides:
            string = ""
            image = []
            for j in i.shapes:
                """#Check if the shape is an image
                try:
                    image.append(j.image)
                except:
                    image.append("")"""
                if not j.has_text_frame:
                    continue
                text_frame = j.text_frame
                for k in text_frame.paragraphs:
                    string += k.text + "\n"
            info.append([string])
            #images.append(image)
        return info

    if file.find(".pdf") >= 0:
        print("PDF Read.")
        prs = PyPDF2.PdfReader(file)
        info = []
        images = []
        for i in prs.pages:
            info.append(i.extract_text())

        return info

        """"#Get Images
        prs = fitz.open(file)
        for i in range(len(prs)):
            page = prs[i]
            images = page.get_images()

            if images:
                for i in images:
                    img = prs.extract_image(i[0])
                    print(img)"""
